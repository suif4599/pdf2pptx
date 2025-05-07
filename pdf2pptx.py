import argparse
import os
import fitz
import warnings
import fractions
import io
import re
import tempfile
import subprocess
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE
from tqdm import tqdm


def extract_link(page: fitz.Page):
    links_info = []
    for link in page.get_links():
        kind = link.get('kind')
        rect = link['from']
        if kind == fitz.LINK_URI:
            uri = link['uri']
            links_info.append((
                "uri",
                uri,
                (
                    rect.x0 / page.rect.width,
                    rect.y0 / page.rect.height,
                    rect.x1 / page.rect.width,
                    rect.y1 / page.rect.height
                ),
            ))
        elif kind == fitz.LINK_GOTO:
            # many times the kind will be LINK_NAMED
            target_page = link.get('page', -1)
            if target_page >= 0:
                links_info.append((
                    "goto",
                    target_page,
                    (
                        rect.x0 / page.rect.width,
                        rect.y0 / page.rect.height,
                        rect.x1 / page.rect.width,
                        rect.y1 / page.rect.height
                    ),
                ))
        elif kind == fitz.LINK_GOTOR:
            file = link.get('file', '')
            target_page = link.get('page', -1)
            uri = file
            if target_page >= 0:
                uri += f"#page={target_page + 1}"
            links_info.append((
                "gotor",
                uri,
                (
                    rect.x0 / page.rect.width,
                    rect.y0 / page.rect.height,
                    rect.x1 / page.rect.width,
                    rect.y1 / page.rect.height
                ),
            ))
        elif kind == fitz.LINK_NAMED:
            target_page = link.get('page', -1)
            if target_page >= 0 and target_page != page.number:
                links_info.append((
                    "goto",
                    target_page,
                    (
                        rect.x0 / page.rect.width,
                        rect.y0 / page.rect.height,
                        rect.x1 / page.rect.width,
                        rect.y1 / page.rect.height
                    ),
                ))
            pass
        else:
            warnings.warn(
                f"Unsupported link kind: {kind}, in page {page.number}"
            )
    return links_info


def add_hyperlink(slide, _type, uri, rect, prs):
    left = int(rect[0] * prs.slide_width)
    top = int(rect[1] * prs.slide_height)
    width = int((rect[2] - rect[0]) * prs.slide_width)
    height = int((rect[3] - rect[1]) * prs.slide_height)
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        left, top, width, height
    )
    if _type == "uri":
        shape.click_action.hyperlink.address = uri
    elif _type == "goto":
        shape.click_action.target_slide = prs.slides[uri]
    elif _type == "gotor":
        shape.click_action.hyperlink.address = uri
    else:
        raise ValueError(f"Unsupported link type: {_type}")
    shape.fill.background()
    shape.line.width = 0
    shape.line.fill.background()
    shape.fill.transparency = 1.0


if __name__ == "__main__":
    parser = argparse.ArgumentParser(
        description='Convert a pdf file to a PowerPoint presentation'
    )
    parser.add_argument(
        'input',
        type=str,
        help='The input pdf file'
    )
    parser.add_argument(
        'output',
        nargs="?",
        default="",
        type=str,
        help='The output PowerPoint file'
    )
    parser.add_argument(
        "-svg",
        action="store_true",
        help="Use SVG for rendering instead of PNG. "
             "This is experimental and depends on inkscape."
    )
    parser.add_argument(
        '--dpi',
        type=int,
        default=600,
        help='The DPI to render the pdf, default is 600. '
             'When using SVG, this is ignored.'
    )
    parser.add_argument(
        '--aspect-ratio',
        type=str,
        default="auto",
        help='The aspect ratio of the slides like "16:9" for example, '
             'default is automatically derived from the pdf'
    )
    parser.add_argument(
        "--inkscape-path",
        type=str,
        default="inkscape",
        help="Path to the inkscape executable. "
             "This is only used when --svg is specified."
    )
    args = parser.parse_args()

    if not os.path.exists(args.input):
        print("Input file does not exist")
        exit(1)

    if not args.output:
        args.output = os.path.splitext(args.input)[0] + ".pptx"

    print(f"Converting {args.input} to {args.output}")

    if args.svg:
        inkscape_path = args.inkscape_path
        proc = subprocess.Popen(
            [inkscape_path, "--version"],
            stdout=subprocess.PIPE,
            stderr=subprocess.PIPE
        )
        stdout, stderr = proc.communicate()
        if proc.returncode != 0:
            print(f"Inkscape not found at {inkscape_path}. "
                  "Please install inkscape or specify the correct path.")
            exit(1)
        print(f"Inkscape version: {stdout.decode().strip()}")

    prs = Presentation()
    layout = prs.slide_layouts[0]
    pdf = fitz.open(args.input)
    rect = pdf.load_page(0).rect
    if args.aspect_ratio == "auto":
        aspect_ratio = rect.width / rect.height
    else:
        match = re.match(r"^(\d+):(\d+)$", args.aspect_ratio)
        if match:
            aspect_ratio = int(match.group(1)) / int(match.group(2))
        else:
            try:
                aspect_ratio = float(args.aspect_ratio)
            except ValueError:
                print("Invalid aspect ratio format. "
                      "Use 'width:height' or a float value.")
                exit(1)
        if fractions.Fraction(aspect_ratio).limit_denominator(100) != \
           fractions.Fraction(rect.width / rect.height).limit_denominator(100):
            warnings.warn(
                "The aspect ratio you provided "
                "does not match the pdf's aspect ratio. "
            )
    aspect_ratio = fractions.Fraction(aspect_ratio).limit_denominator(100)
    if aspect_ratio <= 0:
        print("Aspect ratio must be positive")
        exit(1)
    print(f"Aspect ratio: {aspect_ratio.numerator}:{aspect_ratio.denominator}")

    max_width = 9144000
    max_height = 5143500
    optimal_width = min(max_width, aspect_ratio * max_height)
    optimal_height = optimal_width / aspect_ratio
    prs.slide_width = int(optimal_width)
    prs.slide_height = int(optimal_height)

    for _ in range(pdf.page_count):
        prs.slides.add_slide(layout)

    for i in tqdm(range(pdf.page_count)):
        page = pdf.load_page(i)

        if args.svg:
            img = page.get_svg_image()
            with tempfile.NamedTemporaryFile(suffix=".svg") as temp_svg, \
                 tempfile.NamedTemporaryFile(suffix=".emf") as temp_emf:
                temp_svg.write(img.encode("utf-8"))
                temp_svg.flush()
                subprocess.run(
                    [
                        "inkscape",
                        temp_svg.name,
                        "--export-type=emf",
                        "--export-filename",
                        temp_emf.name
                    ],
                    check=True
                )
                slide = prs.slides[i]
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        slide.shapes._spTree.remove(shape._element)
                slide.shapes.add_picture(
                    temp_emf.name,
                    0,
                    0,
                    prs.slide_width,
                    prs.slide_height
                )
        else:
            img = page.get_pixmap(dpi=args.dpi, matrix=fitz.Matrix(2, 2))
            img_bytes = img.tobytes("png")
            img_stream = io.BytesIO(img_bytes)
            slide = prs.slides[i]
            for shape in slide.shapes:
                if shape.has_text_frame:
                    slide.shapes._spTree.remove(shape._element)
            slide.shapes.add_picture(
                img_stream,
                0,
                0,
                prs.slide_width,
                prs.slide_height
            )
            img_stream.close()

        for _type, uri, rect in extract_link(page):
            add_hyperlink(slide, _type, uri, rect, prs)

    prs.save(args.output)
